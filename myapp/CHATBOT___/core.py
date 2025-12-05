import json
import os
import random
import re
import sys
import time
from pathlib import Path
import base64
import io

from dotenv import load_dotenv
import google.generativeai as genai
from PIL import Image
import requests
from bs4 import BeautifulSoup
import pypdf
import docx
import pptx


def _suppress_warnings():
    class HideStderr:
        def __enter__(self):
            self._original = sys.stderr
            sys.stderr = open(os.devnull, "w")

        def __exit__(self, exc_type, exc_val, exc_tb):
            sys.stderr.close()
            sys.stderr = self._original

    return HideStderr


HideStderr = _suppress_warnings()


load_dotenv()
API_KEY = os.getenv("GEMINI_API_KEY")
if not API_KEY:
    raise ValueError("GEMINI_API_KEY missing. Please set it in the .env file.")

genai.configure(api_key=API_KEY)

PRIMARY_MODEL = "gemini-2.5-flash"
FALLBACK_MODEL = "gemini-2.0-flash-lite"
GEN_CFG = {
    "temperature": 0.5,
    "top_p": 0.9,
    "max_output_tokens": 350,
}

MEM_PATH = Path("neuro_memory.json")
memory = {"topics": []}
if MEM_PATH.exists():
    try:
        memory = json.loads(MEM_PATH.read_text())
    except json.JSONDecodeError:
        memory = {"topics": []}


def _persist_memory():
    MEM_PATH.write_text(json.dumps(memory, indent=2))


def add_topic(topic: str) -> None:
    topic = (topic or "")[:60]
    if not topic:
        return
    if topic not in memory["topics"]:
        memory["topics"].append(topic)
    memory["topics"] = memory["topics"][-200:]
    _persist_memory()


def clean(text: str) -> str:
    return re.sub(r"[#*`]+", "", text or "").strip()


def call_model(prompt: str, model: str = PRIMARY_MODEL, cfg: dict | None = None) -> str:
    cfg = cfg or GEN_CFG
    last_error = None

    try:
        with HideStderr():
            gen = genai.GenerativeModel(model, generation_config=cfg)
            result = gen.generate_content(prompt)
        reply = clean(getattr(result, "text", "") or "")
        if reply:
            return reply
    except Exception as exc:  # pragma: no cover - external service
        last_error = exc
        if "429" in str(exc):
            time.sleep(1.2)

    if last_error:
        try:
            with HideStderr():
                gen = genai.GenerativeModel(FALLBACK_MODEL, generation_config=cfg)
                result = gen.generate_content(prompt)
            reply = clean(getattr(result, "text", "") or "")
            if reply:
                return reply
        except Exception as exc:  # pragma: no cover
            last_error = exc

    if last_error:
        try:
            with HideStderr():
                short_cfg = {"max_output_tokens": 120}
                gen = genai.GenerativeModel(FALLBACK_MODEL, generation_config=short_cfg)
                result = gen.generate_content("Short: " + prompt[:200])
            reply = clean(getattr(result, "text", "") or "")
            if reply:
                return reply
        except Exception as exc:  # pragma: no cover
            last_error = exc

    return f"Model unavailable: {last_error or 'unknown error'}"


def adaptive(profile: str, state: str, user_input: str, history: str | None = None) -> str:
    base = f"""
You are the NeuroAdaptive Learning Companion.

Profile: {profile}
State: {state}
User asked: {user_input}

Respond with SHORT, CLEAN TEXT. NO markdown.
"""

    if history:
        base = base.replace("User asked", f"History:\n{history}\n\nUser asked")

    if state == "attention":
        base += """
Give:
1) Short deep explanation
2) One real example
3) Tiny formula (if applicable)
4) Two related topics
"""

    elif state == "drowsiness":
        base += """
Give:
1) Soft simple explanation
2) One story-like analogy
3) A gentle refocus sentence
"""

    else:
        ask_analogy = random.random() < 0.20
        conceptual = any(word in user_input.lower() for word in ["what", "why", "how", "explain", "define", "concept"])

        base += "Give:\n1) Balanced explanation.\n"

        if ask_analogy and conceptual:
            base += "2) Ask user to create an analogy.\n"
        else:
            base += "2) Do NOT ask for analogy.\n"

    return call_model(base)


def summarize_text(text: str, context: str = "") -> str:
    """Generates a summary for a given text using the LLM."""
    prompt = f"Summarize the following text. {context}:\n\n{text}"
    # Use a generation config optimized for summarization
    summary_cfg = {
        "temperature": 0.3,
        "top_p": 0.8,
        "max_output_tokens": 250,
    }
    return call_model(prompt, cfg=summary_cfg)


def extract_text_from_url(url: str) -> str:
    """Extracts textual content from a URL."""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        # Remove script and style elements
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        # Get text
        text = soup.get_text()
        # Break into lines and remove leading/trailing space on each
        lines = (line.strip() for line in text.splitlines())
        # Break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # Drop blank lines
        return '\n'.join(chunk for chunk in chunks if chunk)
    except requests.RequestException as e:
        print(f"Error fetching URL {url}: {e}")
        return f"Error: Could not retrieve content from the URL."


def extract_text_from_file(file_path: Path) -> str:
    """Extracts text from various file types."""
    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            reader = pypdf.PdfReader(file_path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext == ".docx":
            doc = docx.Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs)
        elif ext == ".pptx":
            pres = pptx.Presentation(file_path)
            return "\n".join(
                shape.text
                for slide in pres.slides
                for shape in pres.shapes
                if hasattr(shape, "text")
            )
        elif ext in [".jpg", ".jpeg", ".png", ".webp"]:
            # For images, we'll use the multimodal capabilities of the model
            # so we just return a marker. The calling function will handle the image data.
            return f"[Image file: {file_path.name}]"
        else: # Plain text
            return file_path.read_text(errors="ignore")
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return f"Error: Could not process the file {file_path.name}."

def analyze_image(image_data: bytes, prompt: str) -> str:
    """Analyzes an image using the multimodal model."""
    try:
        image_parts = [{"mime_type": "image/jpeg", "data": image_data}]
        model = genai.GenerativeModel(PRIMARY_MODEL)
        response = model.generate_content([prompt, *image_parts])
        return clean(response.text)
    except Exception as e:
        print(f"Image analysis error: {e}")
        return "Error: Could not analyze the image."
